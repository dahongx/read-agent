from __future__ import annotations

import asyncio
import logging

from openai import OpenAI
from pptx import Presentation

from app.core.config import settings

logger = logging.getLogger(__name__)

_PLACEHOLDER = "此页为图表页，请参考幻灯片内容"

_NARRATION_PROMPT = """你是一位学术演讲助手。请根据以下PPT幻灯片的文字内容，生成一段简洁的口头讲稿（100-200字）。
要求：
- 用口语化的中文表达，适合演讲朗读
- 围绕幻灯片要点展开，不要逐字照读
- 如果内容是标题页，生成一段简短的开场白

幻灯片内容：
{slide_text}

请直接输出讲稿内容，不要加任何前缀说明。"""


def extract_slide_texts(pptx_path: str) -> list[str]:
    """Extract text content from each slide using python-pptx."""
    prs = Presentation(pptx_path)
    texts: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line.strip():
                        parts.append(line.strip())
        texts.append("\n".join(parts))
    return texts


def _call_llm(slide_text: str) -> str:
    """Call LLM synchronously for one slide's narration."""
    if not slide_text.strip():
        return _PLACEHOLDER

    prompt = _NARRATION_PROMPT.format(slide_text=slide_text)
    client = OpenAI(api_key=settings.LLM_API_KEY, base_url=settings.LLM_BASE_URL)
    response = client.chat.completions.create(
        model=settings.LLM_MODEL,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=400,
    )
    return response.choices[0].message.content or _PLACEHOLDER


async def generate_narrations(slide_texts: list[str]) -> list[str]:
    """Generate narrations for all slides in parallel using asyncio."""
    loop = asyncio.get_event_loop()

    async def _narrate(text: str) -> str:
        try:
            return await loop.run_in_executor(None, _call_llm, text)
        except Exception as exc:
            logger.warning("Narration generation failed for slide: %s", exc)
            return _PLACEHOLDER

    results = await asyncio.gather(*[_narrate(t) for t in slide_texts])
    return list(results)
